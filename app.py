from flask import Flask, request, jsonify, redirect, url_for, render_template, send_file
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
import jwt
import datetime
import uuid
import json
from functools import wraps
from google import genai
import os
import PyPDF2
import io
from docx import Document
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract
import textract
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph


GEMINI_API_KEY = "AIzaSyCthF-5gqz4JIgMqlwyRSuV0EVttxRjaNg"
client = genai.Client(api_key=GEMINI_API_KEY)

app = Flask(__name__)
app.config['SECRET_KEY'] = 'your_secret_key'  # In production, store this securely!
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///studymate.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)

# ============================
#         Models
# ============================

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(200), nullable=False)
    flashcards = db.relationship('Flashcard', backref='owner', lazy=True)
    flashcard_sets = db.relationship('FlashcardSet', backref='owner', lazy=True)
    study_guides = db.relationship('StudyGuide', backref='owner', lazy=True)

class Flashcard(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    question = db.Column(db.Text, nullable=False)
    answer = db.Column(db.Text, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.datetime.utcnow)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    # Link flashcards to a set (if applicable)
    flashcard_set_id = db.Column(db.String(36), db.ForeignKey('flashcard_set.id'))

class FlashcardSet(db.Model):
    __tablename__ = 'flashcard_set'  # Explicit table name to match foreign key in Flashcard
    id = db.Column(db.String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    name = db.Column(db.String(100), nullable=False)
    description = db.Column(db.Text)
    created_at = db.Column(db.DateTime, default=datetime.datetime.utcnow)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    # Relationship to access flashcards in this set
    flashcards = db.relationship('Flashcard', backref='flashcard_set', lazy=True)

class StudyGuide(db.Model):
    id = db.Column(db.String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    title = db.Column(db.String(100), nullable=False)
    content = db.Column(db.Text, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.datetime.utcnow)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    has_pdf = db.Column(db.Boolean, default=False)
    pdf_data = db.Column(db.LargeBinary)  # Store PDF data in the database

# ============================
#    Authentication Utils
# ============================
def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = request.cookies.get("token")  # Retrieve token from cookies

        if not token:
            return redirect(url_for('login'))  # Redirect to login if no token

        try:
            data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=["HS256"])
            current_user = User.query.get(data['user_id'])
            if not current_user:
                raise ValueError("User not found")
        except jwt.ExpiredSignatureError:
            return redirect(url_for('login'))  # Redirect on expired token
        except jwt.InvalidTokenError:
            return redirect(url_for('login'))  # Redirect on invalid token
        except Exception:
            return redirect(url_for('login'))

        return f(current_user, *args, **kwargs)

    return decorated

# ============================
#         Routes
# ============================

@app.route('/')
def index():
    token = request.cookies.get("token")
    
    if token:
        try:
            data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=["HS256"])
            current_user = User.query.get(data['user_id'])
            if current_user:
                return redirect(url_for('dashboard'))
        except jwt.ExpiredSignatureError:
            pass  # Expired token, stay on index page
        except jwt.InvalidTokenError:
            pass  # Invalid token, stay on index page
        except Exception:
            pass  # General failure, stay on index page

    return render_template("index.html")

# ----- Signup Endpoint (GET and POST) -----
@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'GET':
        token = request.cookies.get("token")
        if token:
            try:
                data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=["HS256"])
                current_user = User.query.get(data['user_id'])
                if current_user:
                    return redirect(url_for('dashboard'))
            except:
                pass
        return render_template('signup.html')
    
    data = request.get_json()
    if not data or not all(k in data for k in ('username', 'email', 'password')):
        return jsonify({'message': 'Missing fields'}), 400
    hashed_password = generate_password_hash(data['password'], method='pbkdf2:sha256')
    new_user = User(username=data['username'], email=data['email'], password=hashed_password)
    try:
        db.session.add(new_user)
        db.session.commit()
    except Exception as e:
        return jsonify({'message': 'User could not be created. Username or email might already exist.'}), 400
    return jsonify({'message': 'User created successfully!'}), 201

# ----- Login Endpoint (GET and POST) -----
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'GET':
        token = request.cookies.get("token")
        if token:
            try:
                data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=["HS256"])
                current_user = User.query.get(data['user_id'])
                if current_user:
                    return redirect(url_for('dashboard'))
            except:
                pass
        return render_template('login.html')

    data = request.get_json()
    if not data or not all(k in data for k in ('login', 'password')):
        return jsonify({'message': 'Could not verify'}), 401

    user = User.query.filter((User.username == data['login']) | (User.email == data['login'])).first()

    if not user or not check_password_hash(user.password, data['password']):
        return jsonify({'message': 'Invalid credentials'}), 401

    token = jwt.encode(
        {'user_id': user.id, 'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=24)},
        app.config['SECRET_KEY'],
        algorithm="HS256"
    )
    response = jsonify({'token': token})
    response.set_cookie("token", token, httponly=True)
    return response

@app.route('/logout', methods=['POST'])
@token_required
def logout(current_user):
    response = jsonify({'message': 'Logged out successfully!'})
    response.delete_cookie("token")
    return response, 200

@app.route('/dashboard', methods=['GET'])
@token_required
def dashboard(current_user):
    return render_template('dashboard.html', user=current_user)

@app.route('/get-user-info', methods=['GET'])
@token_required
def get_user_info(current_user):
    return jsonify({'username': current_user.username})

# ============================
#     Flashcard Endpoints
# ============================
@app.route('/flashcards', methods=['GET'])
@token_required
def get_flashcards(current_user):
    flashcard_sets = FlashcardSet.query.filter_by(user_id=current_user.id).all()
    flashcards = Flashcard.query.filter_by(user_id=current_user.id).all()
    return render_template('flashcards.html', flashcard_sets=flashcard_sets, flashcards=flashcards, user=current_user)

@app.route('/flashcards', methods=['POST'])
@token_required
def create_flashcard(current_user):
    """Creates a flashcard manually."""
    data = request.get_json()

    if not data or not all(k in data for k in ('question', 'answer')):
        return jsonify({'message': 'Missing question or answer'}), 400

    new_card = Flashcard(
        question=data['question'],
        answer=data['answer'],
        user_id=current_user.id,
        source="manual"
    )
    db.session.add(new_card)
    db.session.commit()

    return jsonify({'message': 'Flashcard created!', 'flashcard_id': new_card.id}), 201

@app.route('/flashcards/ai', methods=['POST'])
@token_required
def create_ai_flashcards(current_user):
    """Creates flashcards using Gemini AI from either text content or file upload."""
    content = None
    num_flashcards = None

    # Handle file upload
    if 'file' in request.files:
        file = request.files['file']
        if file.filename == '':
            return jsonify({'message': 'No file selected'}), 400
        
        # Get file extension
        file_ext = os.path.splitext(file.filename)[1].lower()
        
        try:
            file_content = file.read()
            
            # Handle different file types
            if file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.csv', '.log', '.xml', '.yaml', '.yml']:
                # Text files - try different encodings
                encodings = ['utf-8', 'latin-1', 'ascii', 'cp1252']
                content = None
                for encoding in encodings:
                    try:
                        content = file_content.decode(encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                if content is None:
                    return jsonify({'message': 'Could not decode file content. Please ensure it is a text file.'}), 400
                    
            elif file_ext == '.pdf':
                # PDF files
                try:
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    content = ""
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
                except Exception as e:
                    return jsonify({'message': f'Error reading PDF file: {str(e)}'}), 400
                    
            elif file_ext in ['.doc', '.docx']:
                # Word documents
                try:
                    doc = Document(io.BytesIO(file_content))
                    content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                except Exception as e:
                    return jsonify({'message': f'Error reading Word document: {str(e)}'}), 400

            elif file_ext in ['.xlsx', '.xls']:
                # Excel files
                try:
                    df = pd.read_excel(io.BytesIO(file_content))
                    # Convert DataFrame to string representation
                    content = df.to_string(index=False)
                except Exception as e:
                    return jsonify({'message': f'Error reading Excel file: {str(e)}'}), 400

            elif file_ext == '.pptx':
                # PowerPoint files
                try:
                    prs = Presentation(io.BytesIO(file_content))
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                except Exception as e:
                    return jsonify({'message': f'Error reading PowerPoint file: {str(e)}'}), 400

            elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
                # Image files - use OCR
                try:                    
                    # Open the image using PIL
                    image = Image.open(io.BytesIO(file_content))
                    
                    # Convert to RGB if necessary (for PNG with transparency)
                    if image.mode in ('RGBA', 'P'):
                        image = image.convert('RGB')
                    
                    # Perform OCR
                    content = pytesseract.image_to_string(image)
                    
                    if not content.strip():
                        return jsonify({'message': 'Could not extract text from image. The image might not contain readable text.'}), 400
                except Exception as e:
                    return jsonify({'message': f'Error processing image file: {str(e)}'}), 400

            else:
                # Try using textract as a fallback for other file types
                try:
                    content = textract.process(io.BytesIO(file_content), extension=file_ext[1:]).decode('utf-8')
                except:
                    return jsonify({'message': 'Unsupported file type. Please upload a text file, PDF, Word document, Excel file, PowerPoint, or image.'}), 400
                
        except Exception as e:
            return jsonify({'message': f'Error reading file: {str(e)}'}), 400
    # Handle text content
    else:
        data = request.get_json()
        if not data or 'content' not in data:
            return jsonify({'message': 'Content is required for AI-generated flashcards'}), 400
        content = data.get('content')
        num_flashcards = data.get('count')

    if not content:
        return jsonify({'message': 'No content provided'}), 400

    try:
        # If no count specified, ask Gemini to determine appropriate number
        if not num_flashcards:
            count_prompt = f"""Based on this content, determine the appropriate number of flashcards needed to effectively cover the material. 
            Consider the complexity and depth of the content. Return only a number between 1 and 30.

Content:
{content}"""

            count_response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=count_prompt,
            )
            try:
                num_flashcards = int(count_response.text.strip())
                # Ensure the number is within reasonable bounds
                num_flashcards = max(1, min(30, num_flashcards))
            except ValueError:
                num_flashcards = 5  # Default if parsing fails

        # Generate flashcards using Google Gemini AI
        prompt = f"""Generate {num_flashcards} high-quality flashcards based on this content. Return ONLY a JSON array containing objects with 'question' and 'answer' fields.

Content to generate flashcards from:
{content}

The response must be in this EXACT format (no other text, just the JSON array):
[
    {{"question": "What is X?", "answer": "X is Y"}},
    {{"question": "What is A?", "answer": "A is B"}}
]

Do not include any markdown, formatting, or additional text. The response should start with '[' and end with ']'."""

        response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=prompt,
        )
        ai_flashcards = response.text.strip()
        
        # Clean up the response to ensure valid JSON
        if not ai_flashcards.startswith('['):
            ai_flashcards = ai_flashcards[ai_flashcards.find('['):]
        if not ai_flashcards.endswith(']'):
            ai_flashcards = ai_flashcards[:ai_flashcards.rfind(']')+1]

        # Parse AI-generated text into flashcards
        flashcards = []
        try:
            # Try to parse as JSON
            cards_data = json.loads(ai_flashcards)
            if isinstance(cards_data, list):
                for card in cards_data:
                    if isinstance(card, dict) and 'question' in card and 'answer' in card:
                        flashcards.append({
                            'question': card['question'],
                            'answer': card['answer']
                        })
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {str(e)}")
            print(f"Received text: {ai_flashcards}")
            return jsonify({'message': 'Could not parse AI response as valid JSON. Please try again.'}), 500

        if not flashcards:
            return jsonify({'message': 'Could not generate valid flashcards. Try different content.'}), 500

        return jsonify({
            'message': 'AI-generated flashcards created!',
            'flashcards': flashcards,
            'count': len(flashcards)
        }), 201

    except Exception as e:
        return jsonify({'message': f'AI generation failed: {str(e)}'}), 500

@app.route('/flashcards/<int:flashcard_id>', methods=['GET'])
@token_required
def get_flashcard(current_user, flashcard_id):
    card = Flashcard.query.filter_by(id=flashcard_id, user_id=current_user.id).first()
    if not card:
        return jsonify({'message': 'Flashcard not found'}), 404
    card_data = {
        'id': card.id,
        'question': card.question,
        'answer': card.answer,
        'created_at': card.created_at
    }
    return jsonify({'flashcard': card_data})

@app.route('/flashcards/<int:flashcard_id>', methods=['PUT'])
@token_required
def update_flashcard(current_user, flashcard_id):
    card = Flashcard.query.filter_by(id=flashcard_id, user_id=current_user.id).first()
    if not card:
        return jsonify({'message': 'Flashcard not found'}), 404
    data = request.get_json()
    if 'question' in data:
        card.question = data['question']
    if 'answer' in data:
        card.answer = data['answer']
    db.session.commit()
    return jsonify({'message': 'Flashcard updated!'})

@app.route('/flashcards/<int:flashcard_id>', methods=['DELETE'])
@token_required
def delete_flashcard(current_user, flashcard_id):
    card = Flashcard.query.filter_by(id=flashcard_id, user_id=current_user.id).first()
    if not card:
        return jsonify({'message': 'Flashcard not found'}), 404
    db.session.delete(card)
    db.session.commit()
    return jsonify({'message': 'Flashcard deleted!'})

# ============================
#   Flashcard Set Endpoints
# ============================
@app.route('/flashcard-sets', methods=['GET'])
@token_required
def get_flashcard_sets(current_user):
    sets = FlashcardSet.query.filter_by(user_id=current_user.id).all()
    output = []
    for s in sets:
        set_data = {
            'id': s.id,
            'name': s.name,
            'description': s.description,
            'created_at': s.created_at
        }
        output.append(set_data)
    return jsonify({'flashcard_sets': output})

@app.route('/flashcard-sets', methods=['POST'])
@token_required
def create_flashcard_set(current_user):
    data = request.get_json()
    if not data or 'name' not in data:
        return jsonify({'message': 'Missing name field'}), 400
    new_set = FlashcardSet(name=data['name'], description=data.get('description', ''), owner=current_user)
    db.session.add(new_set)
    db.session.commit()
    return jsonify({'message': 'Flashcard set created!', 'set_id': new_set.id}), 201

@app.route('/flashcard-sets/<string:set_id>/flashcards', methods=['GET'])
@token_required
def get_flashcards_in_set(current_user, set_id):
    flashcard_set = FlashcardSet.query.filter_by(id=set_id, user_id=current_user.id).first()
    if not flashcard_set:
        return jsonify({'message': 'Flashcard set not found'}), 404
    flashcards = Flashcard.query.filter_by(user_id=current_user.id, flashcard_set_id=set_id).all()
    return jsonify({'flashcards': [{'id': f.id, 'question': f.question, 'answer': f.answer} for f in flashcards]})

@app.route('/flashcard-sets/<string:set_id>/flashcards', methods=['POST'])
@token_required
def create_flashcard_in_set(current_user, set_id):
    flashcard_set = FlashcardSet.query.filter_by(id=set_id, user_id=current_user.id).first()
    if not flashcard_set:
        return jsonify({'message': 'Flashcard set not found'}), 404
    data = request.get_json()
    if not data or not all(k in data for k in ('question', 'answer')):
        return jsonify({'message': 'Missing question or answer'}), 400
    new_flashcard = Flashcard(question=data['question'], answer=data['answer'], user_id=current_user.id, flashcard_set_id=set_id)
    db.session.add(new_flashcard)
    db.session.commit()
    return jsonify({'message': 'Flashcard created!', 'flashcard_id': new_flashcard.id}), 201

@app.route('/flashcard-sets/<string:set_id>', methods=['GET'])
@token_required
def get_flashcard_set(current_user, set_id):
    s = FlashcardSet.query.filter_by(id=set_id, user_id=current_user.id).first()
    if not s:
        return jsonify({'message': 'Flashcard set not found'}), 404
    set_data = {
        'id': s.id,
        'name': s.name,
        'description': s.description,
        'created_at': s.created_at
    }
    return jsonify({'flashcard_set': set_data})

@app.route('/flashcard-sets/<string:set_id>', methods=['PUT'])
@token_required
def update_flashcard_set(current_user, set_id):
    s = FlashcardSet.query.filter_by(id=set_id, user_id=current_user.id).first()
    if not s:
        return jsonify({'message': 'Flashcard set not found'}), 404
    data = request.get_json()
    if 'name' in data:
        s.name = data['name']
    if 'description' in data:
        s.description = data['description']
    db.session.commit()
    return jsonify({'message': 'Flashcard set updated!'})

@app.route('/flashcard-sets/<string:set_id>', methods=['DELETE'])
@token_required
def delete_flashcard_set(current_user, set_id):
    s = FlashcardSet.query.filter_by(id=set_id, user_id=current_user.id).first()
    if not s:
        return jsonify({'message': 'Flashcard set not found'}), 404
    db.session.delete(s)
    db.session.commit()
    return jsonify({'message': 'Flashcard set deleted!'})

# ============================
#    Study Guide Endpoints
# ============================
@app.route('/study-guides', methods=['GET'])
@token_required
def get_study_guides(current_user):
    guides = StudyGuide.query.filter_by(user_id=current_user.id).all()
    study_guides = []
    for guide in guides:
        guide_data = {
            'id': guide.id,
            'name': guide.title,  # Changed from title to name to match template
            'has_pdf': guide.has_pdf  # Use the actual has_pdf value from the database
        }
        study_guides.append(guide_data)
    return render_template('study_guides.html', study_guides=study_guides, user=current_user)

@app.route('/study-guides', methods=['POST'])
@token_required
def create_study_guide(current_user):
    data = request.get_json()
    if not data or 'title' not in data:
        return jsonify({'message': 'Missing title field'}), 400
    # Use the provided manual input for content (if any)
    new_guide = StudyGuide(
        title=data['title'],
        content=data.get('content', ''),  # Updated to include manual input
        user_id=current_user.id
    )
    db.session.add(new_guide)
    db.session.commit()
    return jsonify({'message': 'Study guide created!', 'guide_id': new_guide.id}), 201

@app.route('/study-guides/ai', methods=['POST'])
@token_required
def create_ai_study_guide(current_user):
    try:
        content = None
        title = request.form.get('title') if request.form else request.json.get('title', 'Untitled Study Guide')
        
        # Handle file upload
        if 'file' in request.files:
            file = request.files['file']
            if file.filename == '':
                return jsonify({'message': 'No file selected'}), 400
            
            # Get file extension
            file_ext = os.path.splitext(file.filename)[1].lower()
            
            try:
                file_content = file.read()
                
                # Handle different file types (reuse your existing file processing logic)
                if file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.csv', '.log', '.xml', '.yaml', '.yml']:
                    encodings = ['utf-8', 'latin-1', 'ascii', 'cp1252']
                    for encoding in encodings:
                        try:
                            content = file_content.decode(encoding)
                            break
                        except UnicodeDecodeError:
                            continue
                    if content is None:
                        return jsonify({'message': 'Could not decode file content'}), 400
                        
                elif file_ext == '.pdf':
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    content = ""
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
                        
                elif file_ext in ['.doc', '.docx']:
                    doc = Document(io.BytesIO(file_content))
                    content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                    
                elif file_ext in ['.xlsx', '.xls']:
                    df = pd.read_excel(io.BytesIO(file_content))
                    content = df.to_string(index=False)
                    
                elif file_ext == '.pptx':
                    prs = Presentation(io.BytesIO(file_content))
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                                
                elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
                    image = Image.open(io.BytesIO(file_content))
                    if image.mode in ('RGBA', 'P'):
                        image = image.convert('RGB')
                    content = pytesseract.image_to_string(image)
                    if not content.strip():
                        return jsonify({'message': 'Could not extract text from image'}), 400
                
                else:
                    try:
                        content = textract.process(io.BytesIO(file_content), extension=file_ext[1:]).decode('utf-8')
                    except:
                        return jsonify({'message': 'Unsupported file type'}), 400
                    
            except Exception as e:
                return jsonify({'message': f'Error reading file: {str(e)}'}), 400
                
        else:
            # Handle text input
            data = request.get_json()
            if not data or 'content' not in data:
                return jsonify({'message': 'No content provided'}), 400
            content = data.get('content')

        if not content:
            return jsonify({'message': 'No content to process'}), 400

        # Generate study guide content using Gemini
        prompt = f"""Create a comprehensive and well-structured study guide from the following content. 
        Follow these guidelines:
        1. Begin with a clear introduction or overview
        2. Break down the content into logical sections with clear headings
        3. Use bullet points for key concepts and important details
        4. Include definitions, examples, and explanations where relevant
        5. Highlight important terms or concepts
        6. End with a summary of key takeaways
        7. Format the content in a clean, organized way with proper spacing

Content to process:
{content}

Please ensure the study guide is educational, easy to follow, and retains all important information from the source content."""

        response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=prompt,
        )
        
        study_guide_content = response.text.strip()
        
        # Generate PDF data
        pdf_data = generate_pdf_from_content(study_guide_content)
        
        # Create new study guide with the AI-generated content
        new_guide = StudyGuide(
            title=title,
            content=study_guide_content,
            user_id=current_user.id,
            has_pdf=True,
            pdf_data=pdf_data
        )
        db.session.add(new_guide)
        db.session.commit()
        
        return jsonify({
            'message': 'Study guide created!',
            'guide_id': new_guide.id,
            'has_pdf': True
        }), 201

    except Exception as e:
        db.session.rollback()
        return jsonify({'message': f'Failed to create study guide: {str(e)}'}), 500

@app.route('/study-guides/<string:guide_id>', methods=['GET'])
@token_required
def get_study_guide(current_user, guide_id):
    guide = StudyGuide.query.filter_by(id=guide_id, user_id=current_user.id).first()
    if not guide:
        return jsonify({'message': 'Study guide not found'}), 404
    
    # Check if PDF needs to be regenerated
    if not guide.pdf_data and guide.content:
        try:
            pdf_data = generate_pdf_from_content(guide.content)
            guide.pdf_data = pdf_data
            guide.has_pdf = True
            db.session.commit()
        except Exception:
            guide.has_pdf = False
            db.session.commit()
    
    guide_data = {
        'id': guide.id,
        'title': guide.title,
        'content': guide.content,
        'created_at': guide.created_at,
        'has_pdf': guide.has_pdf
    }
    return jsonify({'study_guide': guide_data})

@app.route('/study-guides/<string:guide_id>', methods=['PUT'])
@token_required
def update_study_guide(current_user, guide_id):
    guide = StudyGuide.query.filter_by(id=guide_id, user_id=current_user.id).first()
    if not guide:
        return jsonify({'message': 'Study guide not found'}), 404
    data = request.get_json()
    if 'title' in data:
        guide.title = data['title']
    if 'content' in data:
        guide.content = data['content']
        # Generate PDF data
        pdf_data = generate_pdf_from_content(data['content'])
        guide.pdf_data = pdf_data
        guide.has_pdf = True
    db.session.commit()
    return jsonify({'message': 'Study guide updated!'})

@app.route('/study-guides/<string:guide_id>', methods=['DELETE'])
@token_required
def delete_study_guide(current_user, guide_id):
    guide = StudyGuide.query.filter_by(id=guide_id, user_id=current_user.id).first()
    if not guide:
        return jsonify({'message': 'Study guide not found'}), 404
    db.session.delete(guide)
    db.session.commit()
    return jsonify({'message': 'Study guide deleted!'})

@app.route('/api/generate-flashcards', methods=['POST'])
@token_required
def generate_flashcards(current_user):
    """Generate flashcards using Gemini AI from text content."""
    data = request.get_json()
    content = data.get('content')
    num_flashcards = data.get('count', 5)

    if not content:
        return jsonify({'message': 'Content is required for AI-generated flashcards'}), 400

    try:
        # Generate flashcards using Google Gemini AI
        prompt = f"""Generate {num_flashcards} high-quality flashcards based on this content:
        
{content}

Return an array of JSON objects in this exact format:
[
    {{"question": "What is the capital of France?", "answer": "Paris"}},
    {{"question": "What is 2+2?", "answer": "4"}}
]

Make the questions clear and concise, and the answers detailed but focused. Ensure the content is accurate and educational. Do not include any additional text or formatting, only the JSON array."""

        response = genai.generate_text(prompt)
        ai_flashcards = response.text

        # Parse AI-generated text into flashcards
        flashcards = []
        try:
            # Try to parse as JSON first
            cards_data = json.loads(ai_flashcards)
            for card in cards_data:
                flashcards.append({
                    'question': card['question'],
                    'answer': card['answer']
                })
        except json.JSONDecodeError:
            # Fallback to line-by-line parsing
            for line in ai_flashcards.strip().split("\n"):
                if ":" in line:
                    question, answer = line.split(":", 1)
                    flashcards.append({
                        'question': question.strip(),
                        'answer': answer.strip()
                    })

        if not flashcards:
            return jsonify({'message': 'Could not generate valid flashcards. Try different content.'}), 500

        return jsonify({
            'message': 'Flashcards generated successfully',
            'flashcards': flashcards
        }), 200

    except Exception as e:
        return jsonify({'message': f'AI generation failed: {str(e)}'}), 500

@app.route('/api/process-file', methods=['POST'])
@token_required
def process_file(current_user):
    """Process uploaded files and generate flashcards using Gemini AI."""
    if 'file' not in request.files:
        return jsonify({'message': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'message': 'No file selected'}), 400

    try:
        # Read file content
        content = file.read().decode('utf-8')
        
        # Generate flashcards using the file content
        num_flashcards = request.form.get('count', 5)
        prompt = f"""Generate {num_flashcards} high-quality flashcards based on this content:
        
{content}

Return an array of JSON objects in this exact format:
[
    {{"question": "What is the capital of France?", "answer": "Paris"}},
    {{"question": "What is 2+2?", "answer": "4"}}
]

Make the questions clear and concise, and the answers detailed but focused. Ensure the content is accurate and educational. Do not include any additional text or formatting, only the JSON array."""

        response = genai.generate_text(prompt)
        ai_flashcards = response.text

        # Parse AI-generated text into flashcards
        flashcards = []
        try:
            # Try to parse as JSON first
            cards_data = json.loads(ai_flashcards)
            for card in cards_data:
                flashcards.append({
                    'question': card['question'],
                    'answer': card['answer']
                })
        except json.JSONDecodeError:
            # Fallback to line-by-line parsing
            for line in ai_flashcards.strip().split("\n"):
                if ":" in line:
                    question, answer = line.split(":", 1)
                    flashcards.append({
                        'question': question.strip(),
                        'answer': answer.strip()
                    })

        if not flashcards:
            return jsonify({'message': 'Could not generate valid flashcards from file.'}), 500

        return jsonify({
            'message': 'Flashcards generated successfully',
            'flashcards': flashcards
        }), 200

    except Exception as e:
        return jsonify({'message': f'File processing failed: {str(e)}'}), 500

def generate_pdf_from_content(content):
    """Helper function to generate a PDF from text content with proper styling."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
    
    # Create styles
    styles = getSampleStyleSheet()
    
    # Custom styles for different elements
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Title'],
        fontSize=24,
        spaceAfter=30,
        leading=30
    )
    
    heading1_style = ParagraphStyle(
        'CustomHeading1',
        parent=styles['Heading1'],
        fontSize=18,
        spaceAfter=16,
        spaceBefore=16,
        leading=22
    )
    
    heading2_style = ParagraphStyle(
        'CustomHeading2',
        parent=styles['Heading2'],
        fontSize=16,
        spaceAfter=14,
        spaceBefore=14,
        leading=20
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=12,
        leading=14
    )
    
    bullet_style = ParagraphStyle(
        'CustomBullet',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=12,
        leading=14,
        leftIndent=20,
        bulletIndent=10
    )
    
    # Process content and convert markdown to styled paragraphs
    story = []
    lines = content.split('\n')
    in_list = False
    
    for line in lines:
        line = line.strip()
        if not line:  # Skip empty lines but add some space
            story.append(Paragraph("<br/>", normal_style))
            continue
            
        # Remove common markdown characters that shouldn't appear in the final text
        line = line.replace('`', '')  # Remove code backticks
        line = line.replace('*', '')  # Remove emphasis asterisks
        line = line.replace('_', '')  # Remove emphasis underscores
        
        # Handle headings
        if line.startswith('# '):
            text = line[2:].strip()
            story.append(Paragraph(text, title_style))
        elif line.startswith('## '):
            text = line[3:].strip()
            story.append(Paragraph(text, heading1_style))
        elif line.startswith('### '):
            text = line[4:].strip()
            story.append(Paragraph(text, heading2_style))
        # Handle bullet points
        elif line.startswith(('- ', '• ', '* ')):
            text = line[2:].strip()
            story.append(Paragraph(f"• {text}", bullet_style))
            in_list = True
        # Handle numbered lists
        elif line[0].isdigit() and line[1:].startswith('. '):
            text = line[line.find('.')+2:].strip()
            number = line[:line.find('.')]
            story.append(Paragraph(f"{number}. {text}", bullet_style))
            in_list = True
        # Regular paragraph
        else:
            if in_list:
                story.append(Paragraph("<br/>", normal_style))
                in_list = False
            story.append(Paragraph(line, normal_style))
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()  # Return the binary PDF data

@app.route('/study-guides/create', methods=['POST'])
@token_required
def create_study_guide_with_pdf(current_user):
    data = request.get_json()
    if not data or 'content' not in data:
        return jsonify({'message': 'Missing content'}), 400
    
    try:
        # Generate a title from the first line or first few words
        content = data['content'].strip()
        title = content.split('\n')[0][:30] if '\n' in content else content[:30]
        
        # Create new study guide
        new_guide = StudyGuide(
            title=title,
            content=content,
            user_id=current_user.id,
            has_pdf=True  # We'll generate the PDF right away
        )
        db.session.add(new_guide)
        db.session.commit()
        
        # Generate and save PDF
        pdf_buffer = generate_pdf_from_content(content)
        pdf_path = os.path.join('static', 'pdfs', f'{new_guide.id}.pdf')
        os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
        
        with open(pdf_path, 'wb') as f:
            f.write(pdf_buffer)
        
        return jsonify({
            'message': 'Study guide created!',
            'guide_id': new_guide.id,
            'has_pdf': True
        }), 201
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'message': f'Failed to create study guide: {str(e)}'}), 500

@app.route('/study-guides/<string:guide_id>/pdf', methods=['GET', 'DELETE'])
@token_required
def handle_study_guide_pdf(current_user, guide_id):
    guide = StudyGuide.query.filter_by(id=guide_id, user_id=current_user.id).first()
    if not guide:
        return jsonify({'message': 'Study guide not found'}), 404
    
    if request.method == 'DELETE':
        try:
            guide.content = ""
            db.session.commit()
            return jsonify({'message': 'PDF deleted!'}), 200
        except Exception as e:
            db.session.rollback()
            return jsonify({'message': f'Failed to delete PDF: {str(e)}'}), 500
    
    # GET method
    if not guide.pdf_data:
        try:
            pdf_data = generate_pdf_from_content(guide.content)
            guide.pdf_data = pdf_data
            guide.has_pdf = True
            db.session.commit()
        except Exception as e:
            return jsonify({'message': f'Failed to generate PDF: {str(e)}'}), 500
    
    return send_file(
        io.BytesIO(guide.pdf_data),
        mimetype='application/pdf',
        as_attachment=False,
        download_name=f'study_guide_{guide.id}.pdf'
    )

# ============================
#         Main
# ============================
if __name__ == '__main__':
    with app.app_context():
        db.create_all()  # Create tables if they don't exist. For production, use migrations.
    app.run(debug=False)
